"""Parser and chunking unit tests without external model execution(P3a 搬运适配版)。

两条依赖 TF 私有样本语料(docx 模板 / 旧版二进制 .xls)的用例改为 skip:
样本属于源仓库的业务文档,不随 SDK 迁移。需要回归真实语料时,把样本放到
本仓库 `docs/samples/` 下并去掉 skip 即可(解析逻辑本身原样搬运,未改动)。
"""

import io
from pathlib import Path
from uuid import UUID

import pytest
from docling.datamodel.base_models import ConversionStatus

from nicekit.kb.image_refs import (
    ImageMarkdownRef,
    parsed_image_occurrences,
    parsed_image_placeholder,
    resolve_parsed_image_placeholders,
)
from nicekit.kb.parsers import ParsedImageCandidate, get_parser, parse_document
from nicekit.kb.parsers.docling import DoclingParser, DoclingParserError
from nicekit.kb.parsing import chunk_text, extract_text


def _png_bytes() -> bytes:
    from PIL import Image

    output = io.BytesIO()
    Image.new("RGB", (40, 30), "red").save(output, format="PNG")
    return output.getvalue()


def _docx_bytes() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("年度参考成本", level=1)
    doc.add_heading("备件", level=2)
    doc.add_paragraph("以下为 2026 参考价。")
    table = doc.add_table(rows=2, cols=3)
    table.rows[0].cells[0].text = "名称"
    table.rows[0].cells[1].text = "项目"
    table.rows[0].cells[2].text = "价格"
    table.rows[1].cells[0].text = "一号机"
    table.rows[1].cells[1].text = "备件|工时"  # 单元格内竖线必须转义
    table.rows[1].cells[2].text = "17元"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_fast_docx_headings_and_gfm_table() -> None:
    md = parse_document(_docx_bytes(), "cost.docx", backend="fast").markdown
    assert "# 年度参考成本" in md
    assert "## 备件" in md
    assert "以下为 2026 参考价。" in md
    assert "| 名称 | 项目 | 价格 |" in md
    assert "|" + " --- |" * 3 in md
    assert "| 一号机 | 备件\\|工时 | 17元 |" in md


def test_fast_docx_extracts_distinct_image_occurrences_in_document_order() -> None:
    from docx import Document
    from PIL import Image

    image_bytes = _png_bytes()
    doc = Document()
    doc.add_paragraph("第一张图片之前")
    doc.add_picture(io.BytesIO(image_bytes))
    doc.add_paragraph("两张图片之间")
    doc.add_picture(io.BytesIO(image_bytes))
    doc.add_paragraph("第二张图片之后")
    output = io.BytesIO()
    doc.save(output)

    parsed = parse_document(output.getvalue(), "gallery.docx", backend="fast")

    assert len(parsed.images) == 2
    first, second = parsed.images
    assert first.source_occurrence != second.source_occurrence
    assert (first.reading_order, second.reading_order) == (0, 1)
    assert first.page is first.slide is first.bbox is None
    assert second.page is second.slide is second.bbox is None
    assert "第一张图片之前" in (first.surrounding_text or "")
    assert "两张图片之间" in (second.surrounding_text or "")
    assert Image.open(io.BytesIO(first.content)).size == (40, 30)
    assert first.content == second.content


@pytest.mark.skip(reason="依赖 TF 私有样本语料(docx 模板),不随 SDK 迁移")
def test_fast_docx_extracts_manifest_header_image_once_from_real_part() -> None:
    sample = (
        Path(__file__).resolve().parents[2]
        / "docs"
        / "20260705-各阶段输出产物模板"
        / "docx"
        / "quote_template.docx"
    )

    parsed = parse_document(sample.read_bytes(), sample.name, backend="fast")

    assert len(parsed.images) == 1
    candidate = parsed.images[0]
    assert candidate.source_occurrence.startswith("docx:word/header1.xml#blip:0:")
    assert candidate.parser_item_ref.startswith("word/header1.xml#blip:0:")
    assert candidate.page is candidate.slide is candidate.bbox is None


def test_fast_xlsx_sheet_heading_and_table() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "设备"
    ws.append(["厂区", "设备", "规格", "参考价"])
    ws.append([None, None, None, None])  # 全空行应跳过
    ws.append(["华东", "Holiday Inn", "4星", 95])
    buf = io.BytesIO()
    wb.save(buf)

    md = parse_document(buf.getvalue(), "devices.xlsx", backend="fast").markdown
    assert "## 工作表:设备" in md
    assert "| 厂区 | 设备 | 规格 | 参考价 |" in md
    assert "| 华东 | Holiday Inn | 4星 | 95 |" in md
    assert "|  |  |  |  |" not in md  # 全空行没混进表格


def test_fast_pptx_preserves_slide_text_tables_and_page_anchors() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "产品概览"
    table = first.shapes.add_table(
        2, 2, Inches(1), Inches(2), Inches(6), Inches(1)
    ).table
    table.cell(0, 0).text = "厂区"
    table.cell(0, 1).text = "数量"
    table.cell(1, 0).text = "华南"
    table.cell(1, 1).text = "3"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "报价说明"
    output = io.BytesIO()
    presentation.save(output)

    parsed = parse_document(output.getvalue(), "product.pptx", backend="fast")

    assert parsed.page_count == 2
    assert parsed.markdown.index("<!--page:1-->") < parsed.markdown.index("产品概览")
    assert "| 厂区 | 数量 |" in parsed.markdown
    assert "| 华南 | 3 |" in parsed.markdown
    assert parsed.markdown.index("<!--page:2-->") < parsed.markdown.index("报价说明")
    assert parsed.parser_name == "fast"


def test_fast_pptx_extracts_picture_shape_with_slide_bbox_and_context() -> None:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(3), Inches(0.5)
    )
    text_box.text = "美泉宫外观"
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes()),
        Inches(1),
        Inches(1),
        width=Inches(2),
        height=Inches(1.5),
    )
    output = io.BytesIO()
    presentation.save(output)

    parsed = parse_document(output.getvalue(), "palace.pptx", backend="fast")

    assert len(parsed.images) == 1
    candidate = parsed.images[0]
    assert candidate.slide == 1 and candidate.page is None
    assert candidate.reading_order == 1
    assert candidate.bbox == pytest.approx(
        (
            1 / (presentation.slide_width / Inches(1)),
            1 / (presentation.slide_height / Inches(1)),
            3 / (presentation.slide_width / Inches(1)),
            2.5 / (presentation.slide_height / Inches(1)),
        )
    )
    assert candidate.surrounding_text == "美泉宫外观"
    assert Image.open(io.BytesIO(candidate.content)).size == (40, 30)


@pytest.mark.skip(reason="依赖 TF 私有样本语料(旧版二进制 .xls),不随 SDK 迁移")
def test_fast_legacy_xls_uses_xlrd_binary_reader() -> None:
    """旧版二进制 .xls 不是 zip 容器,必须走 xlrd 而非 openpyxl。"""
    sample = (
        Path(__file__).resolve().parents[2]
        / "docs"
        / "20260702-客户知识文档部分示例样本"
        / "示例样本目录"
        / "MU3设备list （3选1）- 3.30更新.xls"
    )
    md = parse_document(sample.read_bytes(), sample.name, backend="fast").markdown
    assert "## 工作表:" in md
    assert "| 日期 | 站点 |" in md
    assert "罗马" in md


def test_fast_pdf_page_anchors() -> None:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)

    parsed = parse_document(buf.getvalue(), "brochure.pdf", backend="fast")
    assert "<!--page:1-->" in parsed.markdown
    assert "<!--page:2-->" in parsed.markdown
    assert parsed.page_count == 2


def test_fast_pdf_extracts_embedded_scan_image() -> None:
    from PIL import Image

    source = Image.new("RGB", (80, 60), "blue")
    output = io.BytesIO()
    source.save(output, format="PDF")

    parsed = parse_document(output.getvalue(), "scan.pdf", backend="fast")

    assert len(parsed.images) == 1
    candidate = parsed.images[0]
    assert candidate.page == 1 and candidate.slide is None
    assert candidate.bbox is None
    assert candidate.content_type == "image/jpeg"
    assert candidate.parser_item_ref
    assert Image.open(io.BytesIO(candidate.content)).size == (80, 60)


def test_fast_csv_to_gfm_table() -> None:
    data = "厂区,设备,价格\n华东,Holiday Inn,95\n".encode()
    md = parse_document(data, "devices.csv", backend="fast").markdown
    assert "| 厂区 | 设备 | 价格 |" in md
    assert "| 华东 | Holiday Inn | 95 |" in md


def test_fast_txt_md_passthrough() -> None:
    parsed = parse_document("# 已是 markdown\n正文".encode(), "note.md", backend="fast")
    assert parsed.markdown == "# 已是 markdown\n正文"
    assert parsed.parser_name == "fast"
    assert parse_document("A-100 规格 £95".encode(), "a.txt", backend="fast").markdown == (
        "A-100 规格 £95"
    )


def test_unsupported_type_raises_not_falls_back() -> None:
    with pytest.raises(ValueError):
        parse_document(b"PK\x03\x04", "archive.zip", backend="fast")
    # 非 fast 后端遇到不支持类型同样如实上抛(降级无济于事;
    # docling 对 zip 直接委托 fast,不触发模型加载)
    with pytest.raises(ValueError):
        parse_document(b"PK\x03\x04", "archive.zip", backend="docling")


def test_image_produces_markdown_placeholder() -> None:
    # The fast fallback keeps its compatibility placeholder behavior.
    for name in ("poster.jpg", "map.jpeg", "logo.png", "photo.webp"):
        parsed = parse_document(b"\xff\xd8", name, backend="fast")
        assert parsed.markdown == f"![]({name})"
        assert parsed.parser_name == "fast"
        assert parsed.structured_json is None
        assert len(parsed.images) == 1
        assert parsed.images[0].content == b"\xff\xd8"
        assert parsed.images[0].source_occurrence == "source:0"
        assert parsed.images[0].filename == name
        assert parsed.images[0].reading_order == 0


def test_image_candidate_rejects_page_and_slide_together() -> None:
    with pytest.raises(ValueError, match="mutually exclusive"):
        ParsedImageCandidate(
            content=b"image",
            source_occurrence="fixture:0",
            filename="fixture.png",
            content_type="image/png",
            page=1,
            slide=1,
        )


def test_docling_configures_pdf_and_image_standard_pipeline(monkeypatch) -> None:
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import RapidOcrOptions

    captured: dict = {}

    class FakeConverter:
        def __init__(self, **kwargs) -> None:
            captured.update(kwargs)

    monkeypatch.setattr("docling.document_converter.DocumentConverter", FakeConverter)

    parser = DoclingParser()
    converter = parser._get_converter()

    assert isinstance(converter, FakeConverter)
    format_options = captured["format_options"]
    assert set(format_options) == {InputFormat.PDF, InputFormat.IMAGE}
    pdf_options = format_options[InputFormat.PDF]
    image_options = format_options[InputFormat.IMAGE]
    assert pdf_options.pipeline_cls is image_options.pipeline_cls
    assert pdf_options.pipeline_options is image_options.pipeline_options
    pipeline_options = pdf_options.pipeline_options
    assert pipeline_options.do_ocr is True
    assert pipeline_options.do_table_structure is True
    assert pipeline_options.generate_page_images is True
    assert pipeline_options.generate_picture_images is True
    assert pipeline_options.ocr_batch_size == 1
    assert pipeline_options.layout_batch_size == 1
    assert pipeline_options.table_batch_size == 1
    assert pipeline_options.queue_max_size == 1
    assert isinstance(pipeline_options.ocr_options, RapidOcrOptions)
    assert pipeline_options.ocr_options.lang == ["chinese", "english"]
    assert pipeline_options.ocr_options.force_full_page_ocr is False


def test_docling_image_exports_structured_json_and_markdown_from_same_document() -> None:
    calls: list[str] = []
    structured_json = {
        "schema_name": "DoclingDocument",
        "version": "1.10.0",
        "pages": {"1": {"size": {"width": 100, "height": 200}}},
        "tables": [
            {
                "data": {"table_cells": [{"text": "Bangkok"}]},
                "prov": [{"page_no": 1, "bbox": {"l": 1, "t": 2, "r": 3, "b": 4}}],
            }
        ],
    }

    class FakeDocument:
        pages = {1: object()}

        def export_to_dict(self) -> dict:
            calls.append("structured")
            return structured_json

        def export_to_markdown(self) -> str:
            calls.append("markdown")
            return "| City |\n| --- |\n| Bangkok |"

    class FakeConverter:
        def __init__(self) -> None:
            self.source_name: str | None = None

        def convert(self, stream):
            self.source_name = stream.name
            return type(
                "Result",
                (),
                {
                    "document": FakeDocument(),
                    "status": ConversionStatus.SUCCESS,
                    "errors": [],
                },
            )()

    parser = DoclingParser()
    converter = FakeConverter()
    parser._converter = converter

    parsed = parser.parse(b"image-bytes", "scan.png")

    assert converter.source_name == "scan.png"
    assert calls == ["structured", "markdown"]
    assert parsed.parser_name == "docling"
    assert parsed.page_count == 1
    assert parsed.structured_json == structured_json
    assert parsed.structured_json["tables"][0]["prov"][0]["page_no"] == 1
    assert parsed.markdown == "| City |\n| --- |\n| Bangkok |"
    assert len(parsed.images) == 1
    assert parsed.images[0].content == b"image-bytes"
    assert parsed.images[0].source_occurrence == "source:0"


def test_docling_extracts_picture_pixels_provenance_and_nearby_text() -> None:
    from docling_core.types.doc import (
        BoundingBox,
        CoordOrigin,
        DocItemLabel,
        DoclingDocument,
        ImageRef,
        ProvenanceItem,
        Size,
    )
    from PIL import Image

    doc = DoclingDocument(name="brochure.pdf")
    doc.add_page(page_no=1, size=Size(width=100, height=200))
    doc.add_page(page_no=2, size=Size(width=100, height=200))
    text_provenance = ProvenanceItem(
        page_no=1,
        bbox=BoundingBox(
            l=5,
            t=5,
            r=95,
            b=15,
            coord_origin=CoordOrigin.TOPLEFT,
        ),
        charspan=(0, 10),
    )
    other_page_provenance = ProvenanceItem(
        page_no=2,
        bbox=BoundingBox(
            l=5,
            t=5,
            r=95,
            b=15,
            coord_origin=CoordOrigin.TOPLEFT,
        ),
        charspan=(0, 10),
    )
    picture_provenance = ProvenanceItem(
        page_no=1,
        bbox=BoundingBox(
            l=10,
            t=20,
            r=90,
            b=100,
            coord_origin=CoordOrigin.TOPLEFT,
        ),
        charspan=(0, 0),
    )
    doc.add_text(DocItemLabel.TEXT, "图片之前的说明", prov=text_provenance)
    picture = doc.add_picture(
        image=ImageRef.from_pil(Image.new("RGB", (40, 30), "green"), dpi=72),
        prov=picture_provenance,
    )
    doc.add_text(DocItemLabel.TEXT, "第二页内容不应进入上下文", prov=other_page_provenance)
    doc.add_text(DocItemLabel.TEXT, "图片之后的说明", prov=text_provenance)

    class FakeConverter:
        def convert(self, _stream):
            return type(
                "Result",
                (),
                {
                    "document": doc,
                    "status": ConversionStatus.SUCCESS,
                    "errors": [],
                },
            )()

    parser = DoclingParser()
    parser._converter = FakeConverter()
    parsed = parser.parse(b"%PDF-fixture", "brochure.pdf")

    assert len(parsed.images) == 1
    candidate = parsed.images[0]
    assert candidate.page == 1 and candidate.slide is None
    assert candidate.parser_item_ref == str(picture.self_ref)
    assert candidate.bbox == pytest.approx((0.1, 0.1, 0.9, 0.5))
    assert candidate.reading_order == 1
    assert "图片之前的说明" in (candidate.surrounding_text or "")
    assert "图片之后的说明" in (candidate.surrounding_text or "")
    assert "第二页内容不应进入上下文" not in (candidate.surrounding_text or "")
    assert Image.open(io.BytesIO(candidate.content)).size == (40, 30)
    assert "image" not in parsed.structured_json["pictures"][0]


def _docling_parser_for(doc) -> DoclingParser:
    class FakeConverter:
        def convert(self, _stream):
            return type(
                "Result",
                (),
                {
                    "document": doc,
                    "status": ConversionStatus.SUCCESS,
                    "errors": [],
                },
            )()

    parser = DoclingParser()
    parser._converter = FakeConverter()
    return parser


def test_docling_markdown_anchors_pictures_by_identity_not_position() -> None:
    """占位符与插图是身份对应,不是次序对应。

    取不到像素的插图不进候选,却照样被 docling 导出成 `<!-- image -->`:
    按出现次序回填会把第 1 张候选图贴到这个不属于它的占位上。
    """
    from docling_core.types.doc import DocItemLabel, DoclingDocument, ImageRef
    from PIL import Image

    doc = DoclingDocument(name="brochure.pdf")
    doc.add_text(DocItemLabel.TEXT, "开头正文")
    doc.add_picture()  # 无像素:不进候选,占位符仍会被导出
    doc.add_text(DocItemLabel.TEXT, "红图之前")
    red = doc.add_picture(
        image=ImageRef.from_pil(Image.new("RGB", (40, 30), "red"), dpi=72)
    )
    doc.add_text(DocItemLabel.TEXT, "蓝图之前")
    blue = doc.add_picture(
        image=ImageRef.from_pil(Image.new("RGB", (20, 24), "blue"), dpi=72)
    )

    parsed = _docling_parser_for(doc).parse(b"%PDF-fixture", "brochure.pdf")

    occurrences = [candidate.source_occurrence for candidate in parsed.images]
    assert occurrences == [
        f"docling:{red.self_ref}",
        f"docling:{blue.self_ref}",
    ]
    # 正文里的锚点与候选图逐个同序对应
    assert parsed_image_occurrences(parsed.markdown) == tuple(occurrences)
    # 不属于任何候选的插图保留原生占位,不被冒名顶替
    assert parsed.markdown.count("<!-- image -->") == 1
    markdown = parsed.markdown
    assert markdown.index("开头正文") < markdown.index("<!-- image -->")
    assert markdown.index("<!-- image -->") < markdown.index("红图之前")
    assert markdown.index("红图之前") < markdown.index(
        parsed_image_placeholder(occurrences[0])
    )
    assert markdown.index(parsed_image_placeholder(occurrences[0])) < markdown.index(
        "蓝图之前"
    )
    assert markdown.index("蓝图之前") < markdown.index(
        parsed_image_placeholder(occurrences[1])
    )


def test_docling_markdown_without_candidates_matches_official_export() -> None:
    """没有可回指的插图时,导出结果与 docling 官方入口逐字节一致。"""
    from docling_core.types.doc import DocItemLabel, DoclingDocument

    doc = DoclingDocument(name="notes.pdf")
    doc.add_title(text="标题")
    doc.add_text(DocItemLabel.TEXT, "正文段落")
    doc.add_picture()  # 无像素:不进候选,也就没有锚点

    parsed = _docling_parser_for(doc).parse(b"%PDF-fixture", "notes.pdf")

    assert parsed.images == ()
    assert parsed.markdown == doc.export_to_markdown()


def test_parsed_image_anchors_resolve_to_addressable_asset_refs() -> None:
    """锚点按 source_occurrence 换成资产引用;没有资产的锚点整段删除。"""
    first = UUID("00000000-0000-5000-8000-000000000001")
    second = UUID("00000000-0000-5000-8000-000000000002")
    markdown = "\n\n".join(
        (
            "正文一",
            parsed_image_placeholder("docling:#/pictures/0"),
            "正文二",
            parsed_image_placeholder("docling:#/pictures/1"),
            "正文三",
            parsed_image_placeholder("docling:#/pictures/2"),
        )
    )

    resolved = resolve_parsed_image_placeholders(
        markdown,
        {
            "docling:#/pictures/0": ImageMarkdownRef(asset_id=first, page=12),
            "docling:#/pictures/2": ImageMarkdownRef(asset_id=second),
        },
    )

    assert resolved == "\n\n".join(
        (
            "正文一",
            f"![插图 1（第 12 页）](kb-image://{first})",
            "正文二",
            "",
            "正文三",
            f"![插图 2](kb-image://{second})",
        )
    )
    # 占位符独占一行,替换与删除都不改变行数:已落库的证据行锚点保持有效
    assert resolved.count("\n") == markdown.count("\n")
    assert parsed_image_occurrences(resolved) == ()


def test_parsed_image_placeholder_survives_hostile_source_occurrence() -> None:
    """source_occurrence 里的分隔符必须被转义,不能提前闭合注释。"""
    hostile = "docx:word/header1.xml#blip:0:-->rId1 <script>"

    placeholder = parsed_image_placeholder(hostile)

    assert placeholder.count("-->") == 1
    assert parsed_image_occurrences(f"前 {placeholder} 后") == (hostile,)
    with pytest.raises(ValueError, match="must not be blank"):
        parsed_image_placeholder("  ")


def test_docling_emits_full_page_candidate_for_raster_only_pdf() -> None:
    from docling_core.types.doc import DoclingDocument, ImageRef, Size
    from PIL import Image

    source_image = Image.new("RGB", (80, 60), "blue")
    pdf_output = io.BytesIO()
    source_image.save(pdf_output, format="PDF")

    doc = DoclingDocument(name="scan.pdf")
    doc.add_page(
        page_no=1,
        size=Size(width=80, height=60),
        image=ImageRef.from_pil(source_image, dpi=72),
    )

    class FakeConverter:
        def convert(self, _stream):
            return type(
                "Result",
                (),
                {
                    "document": doc,
                    "status": ConversionStatus.SUCCESS,
                    "errors": [],
                },
            )()

    parser = DoclingParser()
    parser._converter = FakeConverter()
    parsed = parser.parse(pdf_output.getvalue(), "scan.pdf")

    assert len(parsed.images) == 1
    candidate = parsed.images[0]
    assert candidate.source_occurrence == "docling:scanned-page:1"
    assert candidate.page == 1
    assert candidate.bbox == (0.0, 0.0, 1.0, 1.0)
    assert Image.open(io.BytesIO(candidate.content)).size == (80, 60)
    assert "image" not in parsed.structured_json["pages"]["1"]


def test_docling_rejects_partial_conversion_before_publishing_artifacts() -> None:
    parser = DoclingParser()
    parser._converter = type(
        "PartialConverter",
        (),
        {
            "convert": lambda _self, _stream: type(
                "PartialResult",
                (),
                {
                    "status": ConversionStatus.PARTIAL_SUCCESS,
                    "errors": [object()],
                    "document": object(),
                },
            )()
        },
    )()

    with pytest.raises(DoclingParserError) as exc_info:
        parser.parse(b"%PDF-partial", "partial.pdf")

    assert exc_info.value.code == "docling_conversion_incomplete"
    assert exc_info.value.status == str(ConversionStatus.PARTIAL_SUCCESS)


def test_removed_mineru_backend_treated_as_unknown_falls_back_to_fast() -> None:
    # mineru 占位后端已移除(2026-07-23):存量配置里的 "mineru" 按未知后端降级 fast
    parsed = parse_document("纯文本".encode(), "a.txt", backend="mineru")
    assert parsed.parser_name == "fast"
    assert parsed.markdown == "纯文本"


def test_unknown_backend_falls_back_to_fast() -> None:
    parsed = parse_document("纯文本".encode(), "a.txt", backend="no-such-backend")
    assert parsed.parser_name == "fast"


def test_get_parser_unknown_raises() -> None:
    with pytest.raises(ValueError):
        get_parser("no-such-backend")


def test_docling_delegates_table_formats_to_fast() -> None:
    # xlsx/csv/txt/md 不走模型,docling 后端直接复用 fast 转换,不需要 docling 可用
    parsed = parse_document("厂区,价格\n华东,95".encode(), "a.csv", backend="docling")
    assert parsed.parser_name == "fast"
    assert "| 厂区 | 价格 |" in parsed.markdown


def test_docling_failure_is_not_silently_downgraded(monkeypatch) -> None:
    class BrokenDoclingParser:
        def parse(self, _data: bytes, _filename: str):
            raise RuntimeError("OCR model unavailable")

    monkeypatch.setattr(
        "nicekit.kb.parsers.get_parser", lambda _name: BrokenDoclingParser()
    )

    with pytest.raises(RuntimeError, match="OCR model unavailable"):
        parse_document(b"%PDF", "scan.pdf", backend="docling")


def test_extract_text_thin_delegate() -> None:
    assert extract_text("a.txt", "A-100 规格 £95".encode()) == "A-100 规格 £95"
    assert "# 年度参考成本" in extract_text("cost.docx", _docx_bytes())


def test_chunk_text_respects_max_and_covers_all() -> None:
    paragraphs = [f"第{i}段:" + "内容" * 100 for i in range(10)]
    chunks = chunk_text("\n".join(paragraphs), max_chars=500, overlap=50)
    assert all(len(c) <= 500 for c in chunks)
    joined = "".join(chunks)
    for i in range(10):
        assert f"第{i}段" in joined


def test_chunk_text_empty() -> None:
    assert chunk_text("   \n  ") == []
