"""fast 后端:自建结构保留转换(零模型依赖,始终可用的兜底)。

docx 标题层级 → `#` 标题、表格 → GFM 表格、xlsx 每 sheet 一节、
pdf/pptx 每页 `<!--page:N-->` 锚。不再拍平为纯文本。
图片(KB-2)产出 `![](<filename>)` 占位——摄入侧把占位换成 object_key,
caption 可用时再追加描述段(见 ingestion / caption.py)。
"""

import csv
import io
import logging
import re
from pathlib import PurePosixPath

from nicekit.kb.parsers.base import (
    DocumentParser,
    ParsedDocument,
    ParsedImageCandidate,
)

_HEADING_RE = re.compile(r"^Heading (\d+)$")
logger = logging.getLogger(__name__)

# 支持直传的图片类型(KB-2):解析仅产出 markdown 图片占位,不做像素级处理
IMAGE_SUFFIXES = (".jpg", ".jpeg", ".png", ".webp")
_IMAGE_CONTENT_TYPES = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".webp": "image/webp",
}
SUPPORTED_DOCUMENT_SUFFIXES = (
    ".docx",
    ".pptx",
    ".xlsx",
    ".xlsm",
    ".xls",
    ".pdf",
    ".csv",
    ".txt",
    ".md",
    *IMAGE_SUFFIXES,
)


def _bounded_context(values: list[str], *, limit: int = 4000) -> str | None:
    text = "\n".join(value.strip() for value in values if value.strip()).strip()
    return text[:limit] or None


def _nearest_text(values: list[str], index: int) -> str | None:
    selected: list[tuple[int, str]] = []
    if values[index].strip():
        selected.append((index, values[index]))
    for distance in range(1, len(values)):
        for candidate_index in (index - distance, index + distance):
            if (
                0 <= candidate_index < len(values)
                and values[candidate_index].strip()
                and all(existing_index != candidate_index for existing_index, _ in selected)
            ):
                selected.append((candidate_index, values[candidate_index]))
                if len(selected) >= 3:
                    return _bounded_context(
                        [value for _, value in sorted(selected, key=lambda item: item[0])]
                    )
    return _bounded_context([value for _, value in sorted(selected, key=lambda item: item[0])])


def _content_type_for_filename(filename: str) -> str | None:
    return _IMAGE_CONTENT_TYPES.get(PurePosixPath(filename.lower()).suffix)


def _md_cell(value: object) -> str:
    """GFM 单元格转义:| 会破坏列结构,换行会破坏行结构。"""
    text = "" if value is None else str(value)
    return text.replace("|", "\\|").replace("\r\n", "<br>").replace("\n", "<br>").strip()


def _md_table(rows: list[list[str]]) -> str:
    """首行为表头的 GFM 表格;列数按最宽行对齐。"""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    padded = [r + [""] * (width - len(r)) for r in rows]
    lines = ["| " + " | ".join(padded[0]) + " |", "|" + " --- |" * width]
    lines += ["| " + " | ".join(r) + " |" for r in padded[1:]]
    return "\n".join(lines)


class FastParser(DocumentParser):
    name = "fast"

    def parse(self, data: bytes, filename: str) -> ParsedDocument:
        suffix = PurePosixPath(filename.lower()).suffix
        if suffix == ".docx":
            return self._docx(data)
        if suffix == ".pptx":
            return self._pptx(data)
        if suffix in (".xlsx", ".xlsm"):
            return self._xlsx(data)
        if suffix == ".xls":
            return self._xls(data)
        if suffix == ".pdf":
            return self._pdf(data)
        if suffix == ".csv":
            return self._csv(data)
        if suffix in (".txt", ".md"):
            return ParsedDocument(
                markdown=data.decode("utf-8", errors="replace"), parser_name=self.name
            )
        if suffix in IMAGE_SUFFIXES:
            # 图片占位:先用 filename(parser 拿不到 object_key),摄入侧统一替换
            return ParsedDocument(
                markdown=f"![]({filename})",
                parser_name=self.name,
                images=(
                    ParsedImageCandidate(
                        content=data,
                        source_occurrence="source:0",
                        filename=filename,
                        content_type=_IMAGE_CONTENT_TYPES[suffix],
                        reading_order=0,
                        parser_item_ref="source",
                    ),
                ),
            )
        raise ValueError(f"不支持的文件类型: {suffix or filename}")

    def _docx(self, data: bytes) -> ParsedDocument:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        doc = Document(io.BytesIO(data))
        parts: list[str] = []
        # iter_inner_content 保持段落与表格的原始顺序
        for block in doc.iter_inner_content():
            if isinstance(block, Paragraph):
                text = block.text.strip()
                if not text:
                    continue
                style_name = (block.style.name or "") if block.style is not None else ""
                m = _HEADING_RE.match(style_name)
                if m:
                    level = min(int(m.group(1)), 6)
                    parts.append(f"{'#' * level} {text}")
                elif style_name == "Title":
                    parts.append(f"# {text}")
                else:
                    parts.append(text)
            elif isinstance(block, Table):
                rows = [
                    [_md_cell(c.text) for c in row.cells]
                    for row in block.rows
                    if any(c.text.strip() for c in row.cells)
                ]
                if rows:
                    parts.append(_md_table(rows))

        images: list[ParsedImageCandidate] = []
        related_parts = list(doc.part.rels.values())
        header_parts = {
            str(relationship.target_part.partname): relationship.target_part
            for relationship in related_parts
            if relationship.reltype == RT.HEADER
        }
        footer_parts = {
            str(relationship.target_part.partname): relationship.target_part
            for relationship in related_parts
            if relationship.reltype == RT.FOOTER
        }
        xml_parts = [
            doc.part,
            *(header_parts[name] for name in sorted(header_parts)),
            *(footer_parts[name] for name in sorted(footer_parts)),
        ]
        for xml_part in xml_parts:
            part_name = str(xml_part.partname).lstrip("/")
            root = xml_part.element
            paragraph_elements = list(root.iter(qn("w:p")))
            paragraph_indexes = {
                paragraph_element: index
                for index, paragraph_element in enumerate(paragraph_elements)
            }
            paragraph_texts = [
                "".join(
                    text_node.text or ""
                    for text_node in paragraph_element.iter(qn("w:t"))
                ).strip()
                for paragraph_element in paragraph_elements
            ]
            for part_occurrence, blip in enumerate(root.iter(qn("a:blip"))):
                relationship_id = blip.get(qn("r:embed"))
                if not relationship_id:
                    continue
                image_part = xml_part.related_parts.get(relationship_id)
                if image_part is None:
                    logger.warning(
                        "DOCX image relationship is missing in %s: %s",
                        part_name,
                        relationship_id,
                    )
                    continue
                content_type = getattr(image_part, "content_type", "")
                if not str(content_type).startswith("image/"):
                    logger.warning(
                        "DOCX relationship %s in %s is not an image: %s",
                        relationship_id,
                        part_name,
                        content_type,
                    )
                    continue
                image_part_name = str(getattr(image_part, "partname", ""))
                image_filename = PurePosixPath(image_part_name).name or (
                    f"document-image-{len(images) + 1}"
                )
                paragraph_element = blip.getparent()
                while (
                    paragraph_element is not None
                    and paragraph_element.tag != qn("w:p")
                ):
                    paragraph_element = paragraph_element.getparent()
                paragraph_index = (
                    paragraph_indexes.get(paragraph_element)
                    if paragraph_element is not None
                    else None
                )
                surrounding_text = (
                    _nearest_text(paragraph_texts, paragraph_index)
                    if paragraph_index is not None
                    else None
                )
                reading_order = len(images)
                reference = (
                    f"{part_name}#blip:{part_occurrence}:{relationship_id}"
                )
                images.append(
                    ParsedImageCandidate(
                        content=image_part.blob,
                        source_occurrence=f"docx:{reference}",
                        filename=image_filename,
                        content_type=str(content_type),
                        reading_order=reading_order,
                        parser_item_ref=reference,
                        surrounding_text=surrounding_text,
                    )
                )
        return ParsedDocument(
            markdown="\n\n".join(parts),
            parser_name=self.name,
            images=tuple(images),
        )

    def _xlsx(self, data: bytes) -> ParsedDocument:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []
        try:
            for ws in wb.worksheets:
                rows = [
                    [_md_cell(v) for v in row]
                    for row in ws.iter_rows(values_only=True)
                    if any(v is not None and str(v).strip() for v in row)
                ]
                section = f"## 工作表:{ws.title}"
                parts.append(f"{section}\n\n{_md_table(rows)}" if rows else section)
        finally:
            wb.close()
        return ParsedDocument(markdown="\n\n".join(parts), parser_name=self.name)

    def _pptx(self, data: bytes) -> ParsedDocument:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(io.BytesIO(data))
        slides: list[str] = []
        images: list[ParsedImageCandidate] = []
        for index, slide in enumerate(presentation.slides, start=1):
            parts = [f"<!--page:{index}-->"]
            contexts: list[tuple[int, str]] = []
            for shape_index, shape in enumerate(slide.shapes):
                shape_parts: list[str] = []
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    if text:
                        parts.append(text)
                        shape_parts.append(text)
                if getattr(shape, "has_table", False):
                    rows = [
                        [_md_cell(cell.text) for cell in row.cells]
                        for row in shape.table.rows
                        if any(cell.text.strip() for cell in row.cells)
                    ]
                    if rows:
                        table_markdown = _md_table(rows)
                        parts.append(table_markdown)
                        shape_parts.append(table_markdown)
                if shape_parts:
                    contexts.append((shape_index, "\n".join(shape_parts)))
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                image = shape.image
                nearest_context = [
                    value
                    for _, value in sorted(
                        contexts,
                        key=lambda item: (abs(item[0] - shape_index), item[0]),
                    )[:3]
                ]
                left = float(shape.left) / float(presentation.slide_width)
                top = float(shape.top) / float(presentation.slide_height)
                right = float(shape.left + shape.width) / float(
                    presentation.slide_width
                )
                bottom = float(shape.top + shape.height) / float(
                    presentation.slide_height
                )
                images.append(
                    ParsedImageCandidate(
                        content=image.blob,
                        source_occurrence=(
                            f"pptx:slide:{index}:shape:{shape_index}:"
                            f"id:{shape.shape_id}"
                        ),
                        filename=image.filename
                        or f"slide-{index}-image-{shape_index + 1}",
                        content_type=image.content_type,
                        slide=index,
                        bbox=(left, top, right, bottom),
                        reading_order=shape_index,
                        parser_item_ref=(
                            f"pptx:slide:{index}:shape:{shape_index}:"
                            f"id:{shape.shape_id}"
                        ),
                        surrounding_text=_bounded_context(nearest_context),
                    )
                )
            slides.append("\n\n".join(parts))
        return ParsedDocument(
            markdown="\n\n".join(slides),
            page_count=len(presentation.slides),
            parser_name=self.name,
            images=tuple(images),
        )

    def _xls(self, data: bytes) -> ParsedDocument:
        # 旧版二进制 Excel:openpyxl 只认 xlsx zip 容器,必须走 xlrd
        import xlrd

        book = xlrd.open_workbook(file_contents=data)
        parts: list[str] = []
        for sheet in book.sheets():
            rows = [
                [_md_cell(cell.value) for cell in sheet.row(row_index)]
                for row_index in range(sheet.nrows)
                if any(
                    cell.value is not None and str(cell.value).strip()
                    for cell in sheet.row(row_index)
                )
            ]
            section = f"## 工作表:{sheet.name}"
            parts.append(f"{section}\n\n{_md_table(rows)}" if rows else section)
        return ParsedDocument(markdown="\n\n".join(parts), parser_name=self.name)

    def _pdf(self, data: bytes) -> ParsedDocument:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        parts: list[str] = []
        images: list[ParsedImageCandidate] = []
        for page_number, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            parts.append(f"<!--page:{page_number}-->\n{page_text}")
            try:
                page_images = list(page.images)
            except Exception as exc:  # noqa: BLE001 - malformed XObject is local
                logger.warning("PDF page %s image extraction failed: %s", page_number, exc)
                continue
            for page_image_index, image_file in enumerate(page_images):
                image_filename = image_file.name or (
                    f"page-{page_number}-image-{page_image_index + 1}.png"
                )
                source_image_filename = image_filename
                image_data = image_file.data
                content_type = _content_type_for_filename(image_filename)
                if content_type is None:
                    pil_image = getattr(image_file, "image", None)
                    pil_format = getattr(pil_image, "format", None)
                    content_type = {
                        "JPEG": "image/jpeg",
                        "PNG": "image/png",
                        "WEBP": "image/webp",
                    }.get(str(pil_format).upper())
                    if content_type is None and pil_image is not None:
                        output = io.BytesIO()
                        pil_image.save(output, format="PNG")
                        image_data = output.getvalue()
                        image_filename = (
                            f"{PurePosixPath(image_filename).stem}.png"
                        )
                        content_type = "image/png"
                if content_type is None:
                    logger.warning(
                        "PDF page %s image %s has unsupported content type",
                        page_number,
                        image_filename,
                    )
                    continue
                reading_order = len(images)
                images.append(
                    ParsedImageCandidate(
                        content=image_data,
                        source_occurrence=(
                            f"pdf:page:{page_number}:image:{page_image_index}:"
                            f"{source_image_filename}"
                        ),
                        filename=image_filename,
                        content_type=content_type,
                        page=page_number,
                        reading_order=reading_order,
                        parser_item_ref=(
                            f"pypdf:page:{page_number}:image:{page_image_index}:"
                            f"{source_image_filename}"
                        ),
                        surrounding_text=page_text[:4000] or None,
                    )
                )
        return ParsedDocument(
            markdown="\n\n".join(parts),
            page_count=len(reader.pages),
            parser_name=self.name,
            images=tuple(images),
        )

    def _csv(self, data: bytes) -> ParsedDocument:
        text = data.decode("utf-8-sig", errors="replace")
        rows = [
            [_md_cell(v) for v in row]
            for row in csv.reader(io.StringIO(text))
            if any(v.strip() for v in row)
        ]
        return ParsedDocument(markdown=_md_table(rows), parser_name=self.name)
